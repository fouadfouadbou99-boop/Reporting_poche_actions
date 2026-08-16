import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px

from io import BytesIO
import tempfile
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches

from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    PageBreak
)

from reportlab.lib.styles import (
    getSampleStyleSheet
)

# ==================================================
# CONFIGURATION
# ==================================================

st.set_page_config(
    page_title="Reporting Comité RPC",
    page_icon="📊",
    layout="wide"
)

st.title("📊 Reporting Comité RPC V4")

# ==================================================
# FONCTIONS
# ==================================================

def create_excel_report(df):

    output = BytesIO()

    with pd.ExcelWriter(
        output,
        engine="xlsxwriter"
    ) as writer:

        df.to_excel(
            writer,
            sheet_name="Reporting",
            index=False
        )

    output.seek(0)

    return output

# --------------------------------------------------

def save_chart(
    portefeuille,
    benchmark
):

    temp_file = tempfile.NamedTemporaryFile(
        suffix=".png",
        delete=False
    )

    plt.figure(figsize=(8,4))

    plt.plot(
        portefeuille,
        label="Portefeuille",
        linewidth=2
    )

    plt.plot(
        benchmark,
        label="Benchmark",
        linewidth=2
    )

    plt.legend()

    plt.grid(True)

    plt.title(
        "Evolution Base 100"
    )

    plt.savefig(
        temp_file.name,
        bbox_inches="tight"
    )

    plt.close()

    return temp_file.name

# --------------------------------------------------

def create_pdf(

    commentaire,
    kpi_df

):

    buffer = BytesIO()

    doc = SimpleDocTemplate(buffer)

    styles = getSampleStyleSheet()

    content = []

    content.append(
        Paragraph(
            "Reporting Comité RPC",
            styles["Title"]
        )
    )

    content.append(
        Spacer(1,20)
    )

    content.append(
        Paragraph(
            "Synthèse Exécutive",
            styles["Heading2"]
        )
    )

    for _, row in kpi_df.iterrows():

        content.append(
            Paragraph(
                f"{row['Indicateur']} : {row['Valeur']}",
                styles["BodyText"]
            )
        )

    content.append(
        PageBreak()
    )

    content.append(
        Paragraph(
            "Analyse",
            styles["Heading2"]
        )
    )

    content.append(
        Paragraph(
            commentaire,
            styles["BodyText"]
        )
    )

    doc.build(content)

    buffer.seek(0)

    return buffer

# --------------------------------------------------

def create_ppt(

    chart_file,

    perf_port,
    perf_indice,
    alpha,
    beta,
    ir,
    te,
    hit_ratio

):

    prs = Presentation()

    # ----------------------------------

    slide1 = prs.slides.add_slide(
        prs.slide_layouts[0]
    )

    slide1.shapes.title.text = (
        "Reporting Comité RPC"
    )

    slide1.placeholders[1].text = (
        f"""
Performance : {perf_port:.2f} %

Benchmark : {perf_indice:.2f} %

Alpha : {alpha:.2f} %
"""
    )

    # ----------------------------------

    slide2 = prs.slides.add_slide(
        prs.slide_layouts[5]
    )

    slide2.shapes.title.text = (
        "Performance Base 100"
    )

    slide2.shapes.add_picture(
        chart_file,
        Inches(0.5),
        Inches(1.2),
        width=Inches(8)
    )

    # ----------------------------------

    slide3 = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide3.shapes.title.text = (
        "Analyse du Risque"
    )

    slide3.placeholders[1].text = (
        f"""
Bêta : {beta:.2f}

Tracking Error : {te:.2f} %

Volatilité maîtrisée
"""
    )

    # ----------------------------------

    slide4 = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide4.shapes.title.text = (
        "Gestion Active"
    )

    slide4.placeholders[1].text = (
        f"""
Alpha : {alpha:.2f} %

Information Ratio : {ir:.2f}

Hit Ratio : {hit_ratio:.2f} %
"""
    )

    # ----------------------------------

    slide5 = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide5.shapes.title.text = (
        "Recommandations"
    )

    slide5.placeholders[1].text = (
        """
• Analyser les sources de sous-performance

• Réviser l’allocation sectorielle

• Renforcer la sélection de titres

• Surveiller l’évolution du risque actif

• Améliorer la génération d’alpha
"""
    )

    ppt_buffer = BytesIO()

    prs.save(
        ppt_buffer
    )

    ppt_buffer.seek(0)

    return ppt_buffer

# ==================================================
# UPLOAD
# ==================================================

uploaded_file = st.file_uploader(
    "Choisir le fichier Excel",
    type=["xlsx"]
)

if uploaded_file:

    try:

        indicateurs = pd.read_excel(
            uploaded_file,
            sheet_name=0
        )

        data = pd.read_excel(
            uploaded_file,
            sheet_name=2
        )

        indicateurs.columns = [
            "Indicateur",
            "Valeur"
        ]

        d = dict(
            zip(
                indicateurs["Indicateur"],
                indicateurs["Valeur"]
            )
        )

        perf_port = d.get(
            "Performance Portefeuille",
            0
        ) * 100

        perf_indice = d.get(
            "Performance Indice",
            0
        ) * 100

        alpha = d.get(
            "Alpha",
            0
        ) * 100

        beta = d.get(
            "Bêta",
            0
        )

        ir = d.get(
            "Information Ratio",
            0
        )

        te = d.get(
            "Tracking Error Annualisé",
            0
        ) * 100

        hit_ratio = d.get(
            "Hit Ratio",
            0
        ) * 100

        vol = d.get(
            "Volatilité Annualisée Portefeuille",
            0
        ) * 100

        # ==================================================
        # KPI
        # ==================================================

        c1,c2,c3,c4,c5 = st.columns(5)

        c1.metric(
            "Performance",
            f"{perf_port:.2f}%"
        )

        c2.metric(
            "Indice",
            f"{perf_indice:.2f}%"
        )

        c3.metric(
            "Alpha",
            f"{alpha:.2f}%"
        )

        c4.metric(
            "Bêta",
            f"{beta:.2f}"
        )

        c5.metric(
            "IR",
            f"{ir:.2f}"
        )

        # ==================================================
        # BASE100
        # ==================================================

        vl_col = None
        indice_col = None

        for col in data.columns:

            if "VL" in str(col):
                vl_col = col

            if "MAISI" in str(col):
                indice_col = col

        portefeuille = data[vl_col]

        benchmark = data[indice_col]

        base100_portefeuille = (
            portefeuille
            / portefeuille.iloc[0]
        ) * 100

        base100_benchmark = (
            benchmark
            / benchmark.iloc[0]
        ) * 100

        fig = go.Figure()

        fig.add_trace(
            go.Scatter(
                y=base100_portefeuille,
                name="Portefeuille"
            )
        )

        fig.add_trace(
            go.Scatter(
                y=base100_benchmark,
                name="Benchmark"
            )
        )

        fig.update_layout(
            title="Evolution Base 100"
        )

        st.plotly_chart(
            fig,
            use_container_width=True
        )

        chart_file = save_chart(
            base100_portefeuille,
            base100_benchmark
        )

        # ==================================================
        # ANALYSE
        # ==================================================

        commentaire = f"""
Le portefeuille affiche une performance
de {perf_port:.2f}% contre
{perf_indice:.2f}% pour le benchmark.

L'alpha ressort à {alpha:.2f}%.

Le bêta est de {beta:.2f}.

Le tracking error s'établit à
{te:.2f}%.

L'information ratio est de
{ir:.2f}.

Le hit ratio est de
{hit_ratio:.2f}%.
"""

        st.info(commentaire)

        # ==================================================
        # EXPORTS
        # ==================================================

        kpi_df = pd.DataFrame({

            "Indicateur":[
                "Performance",
                "Benchmark",
                "Alpha",
                "Beta",
                "Information Ratio",
                "Tracking Error",
                "Hit Ratio"
            ],

            "Valeur":[
                perf_port,
                perf_indice,
                alpha,
                beta,
                ir,
                te,
                hit_ratio
            ]
        })

        st.download_button(
            "📥 Télécharger Excel",
            create_excel_report(kpi_df),
            file_name="Reporting_Comite.xlsx"
        )

        st.download_button(
            "📄 Télécharger PDF",
            create_pdf(
                commentaire,
                kpi_df
            ),
            file_name="Reporting_Comite.pdf"
        )

        st.download_button(
            "📽 Télécharger PowerPoint",
            create_ppt(
                chart_file,
                perf_port,
                perf_indice,
                alpha,
                beta,
                ir,
                te,
                hit_ratio
            ),
            file_name="Reporting_Comite.pptx"
        )

    except Exception as e:

        st.error(
            f"Erreur : {e}"
        )
